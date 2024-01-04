from flask import Flask, request, jsonify
import os
import sys
from io import BytesIO
import boto3
import ffmpeg
from pptx import Presentation
from moviepy.editor import AudioFileClip, ImageClip, VideoClip, TextClip,concatenate_videoclips
from polly_vtt import PollyVTT
import requests
from botocore.exceptions import NoCredentialsError
#import jsonify
import subprocess
from pdf2image import convert_from_path
from botocore.exceptions import ClientError

app = Flask(__name__)

polly_vtt = PollyVTT()

#Environment Vars & Global Vars
AWS_ACCESS_KEY_ID = os.environ.get('AWS_ACCESS_KEY')  
AWS_SECRET_ACCESS_KEY = os.environ.get('AWS_SECRET_KEY')
bucket_name = os.environ.get('BUCKET_NAME')
AWS_REGION = "us-east-1"

with open('/home/ubuntu/.aws/credentials','w') as file:
    file.write("[default]\n")
    file.write(f'aws_access_key_id={AWS_ACCESS_KEY_ID}\n')
    file.write(f'aws_secret_access_key={AWS_SECRET_ACCESS_KEY}\n')

polly_client = boto3.Session(
    aws_access_key_id=AWS_ACCESS_KEY_ID,
    aws_secret_access_key=AWS_SECRET_ACCESS_KEY,
    region_name=AWS_REGION      
).client("polly")

translate = boto3.client(service_name='translate', region_name=AWS_REGION, use_ssl=True)

def extract_slide_notes(prs,speed,LANGUAGE):
    slide_notes = []
    for slide in prs.slides:
        notes_slide = slide.notes_slide
        if notes_slide:
            spokentext = ''
            if LANGUAGE == 'EN':
                spokentext = notes_slide.notes_text_frame.text
            if LANGUAGE == 'ES':
                result = translate.translate_text(Text=notes_slide.notes_text_frame.text, 
                    SourceLanguageCode="en", TargetLanguageCode="es")
                spokentext = result.get('TranslatedText')
            ssmltext = f'<speak><prosody rate="{speed}%">' + spokentext + '</prosody></speak>'
            slide_notes.append(ssmltext)
    return slide_notes


def add_static_image_to_audio(image_path, audio_path, caption_path):

    """Create and save a video file to `output_path` after 
    combining a static image that is located in `image_path` 
    with an audio file in `audio_path`"""
    # create the audio clip object
    audio_clip = AudioFileClip(audio_path)
    # create the image clip object
    image_clip = ImageClip(image_path)
    # use set_audio method from image clip to combine the audio with the image
    video_clip = image_clip.set_audio(audio_clip)
    # specify the duration of the new clip to be the duration of the audio clip
    video_clip.duration = audio_clip.duration
    # set the FPS to 1
    video_clip.fps = 1
    return video_clip

session = boto3.Session()
s3_cli = session.client('s3')

@app.route('/health', methods=['GET'])
def health_check():
    return "", 200

@app.route('/presignedurl',methods=['GET'])
def generate_presignedurl():
    if bucket_name == " " or bucket_name is None:
        raise Exception("Invalid Bucket Name !!")
    
    object_name = request.form.get('object_name')
    #object_name = data.get('object_name')
    if object_name == " " or object_name is None:
        raise Exception("Invalid object name received")

    try:
        presigned_url = create_presigned_url(bucket_name, object_name)
        return jsonify({'presigned_url': presigned_url}), 200 
    except:
        raise Exception("Some Error in Creating in PresignedUrl !!")
        return None, 500

def create_presigned_url(bucket_name, object_name):
    try:
        response = s3_cli.generate_presigned_post(bucket_name, object_name, ExpiresIn=150)
    except NoCredentialsError:
        print("Credentials not available")
        return None
    return response

def create_download_presigned_url(bucket_name, object_name):
    # try:
    #     response = s3_cli.generate_presigned_url(
    #         'get_object',
    #         Params={'Bucket': bucket_name, 'Key': object_name},
    #         ExpiresIn=300
    #     )
    # except NoCredentialsError:
    #     print("Credentials not available")
    #     return None
    # except ClientError as e:
    #     print(e)
    #     return None
    aws_cli_command = ['aws', 's3', 'presign',f's3://{bucket_name}/{object_name}','--expires-in', '300']
    result = subprocess.run(aws_cli_command, capture_output=True, text=True)
    if result.returncode != 0:
        print(f"Error: {result.stderr}")
    return result.stdout.strip()

@app.route('/videolink',methods=['GET'])
def convert_pptx_to_mp4():
    # Parameters collected from the Frontend
    object_name = request.form.get('object_name')
    LANGUAGE = request.form.get('language')
    gender = request.form.get('gender')
    speed = request.form.get('speed')
    if speed=="fast":
        speed = 100
    elif speed=="medium":
        speed = 92
    elif speed=="slow":
        speed = 80

    if object_name is None:
        return "Empty File Name Received", 400
    PPTX_FILE = "pptx/"+object_name
    OUTPUT_VIDEO = "output/"+ object_name.replace('.pptx','')+".mp4"
    
    try:
        s3_client = boto3.client('s3')
        s3_client.download_file(bucket_name, object_name, PPTX_FILE)
        print("File Downloaded successfully.......")
    except Exception as e:
        #print(f"Error downloading the file: {e}")
        return f"Error downloading the file: Please check your credentials: {e}",500
    
    command = ['soffice','--headless','--convert-to','pdf','--outdir','pptx', PPTX_FILE]              #outdir params
    process = subprocess.Popen(command)
    process.wait()

    pdf_file_name = PPTX_FILE.replace('.pptx','.pdf')
    images = convert_from_path(pdf_file_name)

    for i, image in enumerate(images):
        image.save(f"slides/{object_name.split('.')[0]}Slide{i+1}.jpeg", 'JPEG')
    
    
    prs = Presentation(PPTX_FILE)
    slide_notes = extract_slide_notes(prs,speed,LANGUAGE)

    video_clips = []
    for i, slide in enumerate(prs.slides):
        # Convert slide notes to speech
        print(slide.slide_id)
        notes = slide_notes[i]
        #speech_data = convert_text_to_speech(notes)
        file_name = f"{object_name.split('.')[0]}mp3_file{i+1}"
        if LANGUAGE == 'EN':
            # mp3 with VTT captions
            if gender == "male":
                polly_vtt.generate(
                    file_name,
                    VoiceId="Matthew",
                    OutputFormat="mp3",
                    Engine="neural",
                    TextType='ssml',
                    Text=notes
                )
            else:
                polly_vtt.generate(
                    file_name,
                    VoiceId="Joanna",
                    OutputFormat="mp3",
                    Engine="neural",
                    TextType='ssml',
                    Text=notes
                )
        elif LANGUAGE == 'ES':
            # mp3 with VTT captions
            if gender == "male":
                polly_vtt.generate(
                    file_name,
                    VoiceId="Andres",
                    OutputFormat="mp3",
                    Engine="neural",
                    LanguageCode="es-MX",
                    TextType='ssml',
                    Text=notes
                )
            else:
                polly_vtt.generate(
                    file_name,
                    VoiceId="Mia",
                    OutputFormat="mp3",
                    Engine="neural",
                    LanguageCode="es-MX",
                    TextType='ssml',
                    Text=notes
                )

        #imagepath
        image_path = f"slides/{object_name.split('.')[0]}Slide{i+1}.jpeg"
        mp3_filename = "audio-vtt/"+file_name+".mp3"
        vtt_filename = "audio-vtt/"+file_name+".mp3"+".vtt"
        os.rename(file_name+".mp3", mp3_filename)
        os.rename(file_name+".mp3.vtt", vtt_filename)
        individual_video_clip = add_static_image_to_audio(image_path, mp3_filename, vtt_filename)
        mp4_filename = "slide-videos/"+file_name+".mp4"
        individual_video_clip.write_videofile(mp4_filename,fps=1,audio_codec='aac')
        video_clips.append(individual_video_clip)
    final = concatenate_videoclips(video_clips)
    #writing the video into a file / saving the combined video
    final.write_videofile(OUTPUT_VIDEO,fps=1,audio_codec='aac')
    
    try:
        s3_client.upload_file(OUTPUT_VIDEO, 'ppt2video-test', OUTPUT_VIDEO.replace("output/",""))
        print("File uploaded successfully to S3.")
    except Exception as e:
        print(f"Error uploading file: {e}")
    download_url = create_download_presigned_url(bucket_name, OUTPUT_VIDEO.replace("output/",""))
    if download_url.startswith('https'):
        print('presigned_URl: ',download_url)
        cleanup(object_name)
        return jsonify(downloadurl=download_url), 200
    else:
        return jsonify(error='Error in Download link generation'), 500
def cleanup(object_name):
    #clean pdf
    filePrefix = object_name.split('.')[0]
    files = os.listdir('pptx')                          #pptx
    # Iterate over the files and delete them
    for file in files:
        if filePrefix in file:
            os.remove('pptx/'+ file)

    # Get all the files in the given directory
    files = os.listdir('audio-vtt')
    # Iterate over the files and delete them
    for file in files:
        if filePrefix in file:
            os.remove('audio-vtt/'+file)

    # Get all the files in the given directory
    files = os.listdir('slide-videos')
    # Iterate over the files and delete them
    for file in files:
        if filePrefix in file:
            os.remove('slide-videos/'+file)

     # Get all the files in the given directory
    files = os.listdir('slides')
    # Iterate over the files and delete them
    for file in files:
        if filePrefix in file:
            os.remove('slides/'+file)

    # Get all the files in the given directory
    files = os.listdir('output')
    # Iterate over the files and delete them
    for file in files:
        if filePrefix in file:
            os.remove('output/'+file)

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0')
